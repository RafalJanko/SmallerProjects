from pptx import Presentation
import glob
import pandas as pd
from openpyxl import Workbook

book = Workbook()
sheet = book.active
data = []
col = 0
row = 0

# find files in a location and extract all data written as "text object"

for eachfile in glob.glob("PATH TO YOUR FOLDER WITH PPTS"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                new = shape.text
                sep = new.split("\n")
                for item in sep:
                    data.append(sep)

# create a df for pandas to read
df = pd.DataFrame(data)
# create excel writer object
writer = pd.ExcelWriter('output.xlsx')
# write dataframe to excel
df.to_excel(writer)
# save the excel
writer.save()
print('DataFrame is written successfully to Excel File.')
